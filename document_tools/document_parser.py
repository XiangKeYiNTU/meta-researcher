# from camel.toolkits import (
#     ExcelToolkit,
#     AudioAnalysisToolkit,
#     VideoAnalysisToolkit,
#     ImageAnalysisToolkit
# )

import xmltodict
import json
from pptx import Presentation
from docx import Document
from pathlib import Path
# import subprocess
import zipfile
# from pdfminer.pdfparser import PDFParser
# from pdfminer.pdfdocument import PDFDocument
# import pdftotext
import pdfplumber
import os
import pandas as pd
from openai import OpenAI
from dotenv import load_dotenv
import base64

# Extensions that can be directly decoded by UTF-8
UTF8_EXTENSIONS = [".txt", ".md", ".csv", ".json", ".jsonl", "jsonld", ".py"]

# Get the path to the parent folder
parent_env_path = Path(__file__).resolve().parents[1] / ".env"

# Load the .env file from the parent folder
load_dotenv(dotenv_path=parent_env_path)

OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY")

def encode_image(image_path: str) -> str:
    """
    Encodes an image file to a base64 string.
    
    Args:
        image_path (str): Path to the image file.
        
    Returns:
        str: Base64 encoded string of the image.
    """
    with open(image_path, "rb") as img_file:
        return base64.b64encode(img_file.read()).decode('utf-8')

class DocumentParser:
    def __init__(self):
        # self.excel_toolkit = ExcelToolkit()
        # self.audio_toolkit = AudioAnalysisToolkit()
        # self.video_toolkit = VideoAnalysisToolkit()
        # self.image_toolkit = ImageAnalysisToolkit()
        self.temp_dir = Path(__file__).resolve().parents[1] / "temp"
        self.qwen_client = OpenAI(
                                base_url="https://openrouter.ai/api/v1",
                                api_key=OPENROUTER_API_KEY,
                            )


    def _utf8_decode(self, file_path: str):
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()
        return content
    

    def parse_excel_using_pandas(self, file_path: str):
        """
        Parses an Excel file using pandas and returns a DataFrame.
        """
        try:
            df = pd.read_excel(file_path)
            return df.to_string(index=False)  # Convert DataFrame to string for easier readability
        except Exception as e:
            print(f"Error parsing Excel file: {e}")
            return None

    def parse_image_using_qwen(self, image_path: str):
        """
        Uses Qwen model to analyze an image and return a description.
        """
        encoded_image = encode_image(image_path)
        completion = self.qwen_client.chat.completions.create(
            model="qwen/qwen2.5-vl-32b-instruct:free",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "Describe the content of this image in detail."
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/{image_path.split('.')[-1]};base64,{encoded_image}"
                            }
                        }
                    ]
                }
            ]
        )

        return completion.choices[0].message.content
    
    def parse_pdf(self, file_path: str):
        page_contents = []

        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                # Extract plain text
                text = page.extract_text() or ""

                # Extract tables
                tables = page.extract_tables()
                table_texts = []
                for table in tables:
                    # Convert each row into tab-separated text
                    table_rows = ["\t".join(cell if cell is not None else "" for cell in row) for row in table]
                    table_texts.append("\n".join(table_rows))

                combined = f"--- Page {page_num} ---\n"
                combined += f"\n{text.strip()}\n" if text else ""
                if table_texts:
                    combined += "\n\n[Tables:]\n" + "\n\n".join(table_texts)

                page_contents.append(combined.strip())

        return "\n\n".join(page_contents)
    
    def parse_xml(self, file_path: str):
        data = None
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()
        f.close()

        data = xmltodict.parse(content)
        return json.dumps(data, indent=2)

    # def parse_excel(self, file_path):
    #     # return a report of the Excel file content
    #     content = self.excel_toolkit.extract_excel_content(file_path)
    #     return content

    # def analyze_audio(self, file_path):
    #     content = self.audio_toolkit.ask_question_about_audio(audio_path=file_path, question="Explain the content of this audio file as detailed as possible.")
    #     return content

    # def analyze_video(self, file_path):
    #     content = self.video_toolkit.ask_question_about_video(video_path=file_path, question="Explain the content of this video file as detailed as possible.")
    #     return content

    # def parse_xml(self, file_path):
    #     # Implement XML parsing logic here
    #     data = None
    #     with open(file_path, "r", encoding="utf-8") as f:
    #         content = f.read()
    #     f.close()

    #     try:
    #         data = xmltodict.parse(content)
    #         print(f"The extracted xml data is: {data}")
    #         return data

    #     except Exception:
    #         print(f"The raw xml data is: {content}")
    #         return content
        
    # def parse_json(self, file_path):
    #     with open(file_path, "r", encoding="utf-8") as f:
    #         content = json.load(f)
    #     f.close()
    #     return content

    def parse_ppt(self, file_path):
        # Implement PPT parsing logic here
        prs = Presentation(file_path)
        # text_runs = []
        content = []

        for slide_num, slide in enumerate(prs.slides):
            slide_content = f"Content of slide {slide_num}:\n"
            for shape_num, shape in enumerate(slide.shapes):
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            slide_content += run.text
                            slide_content += "\n"
                elif shape.shape_type == 13:  # If the shape is a picture
                    image = shape.image
                    image_bytes = image.blob
                    image_ext = image.ext  # e.g., 'jpeg', 'png'

                    filename = f"slide{slide_num+1}_img{shape_num+1}.{image_ext}"
                    self.temp_dir.mkdir(parents=True, exist_ok=True)
                    image_path = self.temp_dir / filename
                    with open(image_path, "wb") as img_file:
                        img_file.write(image_bytes)

                    image_content = self.parse_image_using_qwen(image_path)
                    slide_content += f"Image {shape_num+1} content: {image_content}\n"
            content.append(slide_content)

        return '\n\n'.join(content)

    def _unzip_file(self, zip_path: str):
        if not zip_path.endswith(".zip"):
            raise ValueError("Only .zip files are supported")

        zip_name = os.path.splitext(os.path.basename(zip_path))[0]
        extract_path = os.path.join(str(self.temp_dir), zip_name)
        os.makedirs(extract_path, exist_ok=True)

        # try:
        #     subprocess.run(["unzip", "-o", zip_path, "-d", extract_path], check=True)
        # except subprocess.CalledProcessError as e:
        #     raise RuntimeError(f"Failed to unzip file: {e}")

        # extracted_files = []
        # for root, _, files in os.walk(extract_path):
        #     for file in files:
        #         extracted_files.append(os.path.join(root, file))

        # return extracted_files

        extracted_files = []

        with zipfile.ZipFile(zip_path, 'r') as zip_ref:
            for member in zip_ref.namelist():
                # Extract the member and get the full path
                zip_ref.extract(member, path=extract_path)
                full_path = os.path.join(extract_path, member)
                extracted_files.append(os.path.normpath(full_path))  # Normalize for Windows

        return extracted_files
    
    def parse_zip(self, file_path: str):
        """
        Parses a zip file and returns the content of each file inside it.
        """
        extracted_files = self._unzip_file(file_path)
        results = []
        for file in extracted_files:
            content = ""
            if file.endswith(tuple(UTF8_EXTENSIONS)):
                content = self._utf8_decode(file)
            elif file.endswith(".xlsx") or file.endswith(".xls") or file.endswith(".xlsm"):
                content = self.parse_excel_using_pandas(file)
            elif file.endswith(".pptx") or file.endswith(".ppt"):
                content = self.parse_ppt(file)
            elif file.endswith(".jpg") or file.endswith(".png") or file.endswith(".jpeg"):
                content = self.parse_image_using_qwen(file)
            elif file.endswith(".pdf"):
                content = self.parse_pdf(file)
            elif file.endswith(".docx") or file.endswith(".doc"):
                content = self.parse_doc(file)
            elif file.endswith(".xml"):
                content = self.parse_xml(file)
            else:
                content = f"Unsupported file type: {file}"
            results.append({os.path.basename(file): content})
        return json.dumps(results, indent=2)

    def parse_doc(self, doc_path: str):
        """
        Parse a Word document (.docx) and extract its text content.
        
        Args:
            doc_path (str): Path to the Word document.
        
        Returns:
            str: Extracted text content from the document.
        """
        try:
            doc = Document(doc_path)
        except Exception as e:
            raise ValueError(f"Error opening document '{doc_path}': {e}")

        paragraphs = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:  # Ignore empty paragraphs
                paragraphs.append(text)

        # Optionally handle tables too
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    paragraphs.append(" | ".join(row_text))

        return "\n".join(paragraphs)

    def parse_file(self, file_path: str):
        """
        Parse a file based on its extension.
        """
        ext = Path(file_path).suffix.lower()
        if ext in UTF8_EXTENSIONS:
            return self._utf8_decode(file_path)
        elif ext in [".xlsx", ".xls", ".xlsm", ".csv"]:
            return self.parse_excel_using_pandas(file_path)
        elif ext in [".pptx", ".ppt"]:
            return self.parse_ppt(file_path)
        elif ext in [".jpg", ".png", ".jpeg"]:
            return self.parse_image_using_qwen(file_path)
        elif ext == ".pdf":
            return self.parse_pdf(file_path)
        elif ext in [".docx", ".doc"]:
            return self.parse_doc(file_path)
        elif ext == ".zip":
            return self.parse_zip(file_path)
        elif ext in [".xml"]:
            return self.parse_xml(file_path)
        else:
            return f"Unsupported file type: {ext}"

    # def parse_doc(self, doc_path: str):

    #     ext = Path(doc_path).suffix.lower()
    #     # if ext in [".xlsx", ".xls", ".xlsm", ".csv"]:
    #     #     return self.parse_excel(doc_path)
    #     # elif ext in [".mp3", ".wav", ".flac"]:
    #     #     return self.analyze_audio(doc_path)
    #     # elif ext in [".mp4", ".avi", ".mov"]:
    #     #     return self.analyze_video(doc_path)
    #     if ext in [".xml"]:
    #         return self.parse_xml(doc_path)
    #     elif ext in [".pptx", ".ppt"]:
    #         return self.parse_ppt(doc_path)
    #     elif ext == ".zip":
    #         extracted_files = self._unzip_file(doc_path)
    #         results = []
    #         for file in extracted_files:
    #             result = self.parse_doc(file)
    #             results.append(result)
    #         return results
    #     elif ext in [".json", ".jsonl"]:
    #         return self.parse_json(doc_path)
    #     else:
    #         with open(doc_path, "r", encoding="utf-8") as f:
    #             content = f.read()
    #         return content
    
if __name__ == "__main__":
    # # Get the path to the parent folder
    # parent_env_path = Path(__file__).resolve().parents[1] / ".env"

    # # Load the .env file from the parent folder
    # load_dotenv(dotenv_path=parent_env_path)

    parser = DocumentParser()
    # Example usage
    file_path = ""  # Change this to your file path
    result = parser.parse_file(file_path)
    print(result)